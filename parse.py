from qdrant_client import QdrantClient, models
from qdrant_client.models import PointStruct
from fastembed.embedding import TextEmbedding
from qdrant_client.models import Filter
from fastembed import SparseTextEmbedding
import tqdm
from fastembed import ImageEmbedding
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from PIL import Image
import numpy as np
import os
import json
import uuid

def extract_content_from_pptx(pptx_path, output_dir):
    presentation = Presentation(pptx_path)
    dataset = []

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for slide_number, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_number": slide_number + 1,
            "text": None,
            "image_paths": []
        }

        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        slide_data["text"] = slide_text.strip() if slide_text.strip() else None

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                image_filename = f"slide_{slide_number + 1}_image_{shape.shape_id}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                slide_data["image_paths"].append(image_path)

        dataset.append(slide_data)

    return dataset


def process_presentations(folder_path, output_dir):
    all_presentations_data = []

    for file_name in os.listdir(folder_path):
        if file_name.endswith(".pptx"):
            pptx_path = os.path.join(folder_path, file_name)
            print(f"Обработка файла: {pptx_path}")
            presentation_data = extract_content_from_pptx(pptx_path, output_dir)
            all_presentations_data.append({
                "file_name": file_name,
                "slides": presentation_data
            })

    return all_presentations_data

presentations_folder = "presentations"  
output_images_folder = "extracted_images"  

dataset = process_presentations(presentations_folder, output_images_folder)