from qdrant_client import QdrantClient, models
from qdrant_client.models import PointStruct
from fastembed.embedding import TextEmbedding
from qdrant_client.models import Filter
from fastembed import SparseTextEmbedding
import tqdm
from fastembed import ImageEmbedding
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from PIL import Image
import numpy as np
import os
import json
import uuid

# URL та API ключ твого кластера
api_key = "SvkHzA7fxPiIwvrsBq0oruGXzKxf4zIm-G8ATUJxb7GMrjvgN6QyfA"
url = "https://88e4a23d-b0b1-4348-b80d-94c77729aa8a.us-east4-0.gcp.cloud.qdrant.io:6333"

# Підключення до кластера
client = QdrantClient(url=url, api_key=api_key)
client.set_sparse_model("Qdrant/bm42-all-minilm-l6-v2-attentions")

# Инициализация модели для текста и изображений
text_model = TextEmbedding("sentence-transformers/all-MiniLM-L6-v2")  # Для текста
text_bm25 = SparseTextEmbedding("Qdrant/bm42-all-minilm-l6-v2-attentions")
image_model = ImageEmbedding(model_name="Qdrant/clip-ViT-B-32-vision")

def extract_content_from_pptx(pptx_path, output_dir):
    presentation = Presentation(pptx_path)
    dataset = []

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for slide_number, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_number": slide_number + 1,
            "text": None,
            "image_paths": []
        }

        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        slide_data["text"] = slide_text.strip() if slide_text.strip() else None

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                image_filename = f"slide_{slide_number + 1}_image_{shape.shape_id}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                slide_data["image_paths"].append(image_path)

        dataset.append(slide_data)

    return dataset


def process_presentations(folder_path, output_dir):
    all_presentations_data = []

    for file_name in os.listdir(folder_path):
        if file_name.endswith(".pptx"):
            pptx_path = os.path.join(folder_path, file_name)
            print(f"Обработка файла: {pptx_path}")
            presentation_data = extract_content_from_pptx(pptx_path, output_dir)
            all_presentations_data.append({
                "file_name": file_name,
                "slides": presentation_data
            })

    return all_presentations_data

presentations_folder = "presentations"  
output_images_folder = "extracted_images"  

dataset = process_presentations(presentations_folder, output_images_folder)


texts = [slide["text"] for presentation in dataset for slide in presentation["slides"]]
images = []
for presentation in dataset:
    for slide in presentation["slides"]:
        for image_path in slide["image_paths"]:
            img = Image.open(image_path)
            images.append(img)


text_model_embeddings = list(text_model.embed(texts))  # Эмбеддинг для текста
image_embeddings = list(image_model.embed(images))  # Эмбеддинг для изображений

batch_size = 16
encoded_sparse_docs = list(client._sparse_embed_documents(    documents=texts,
    embedding_model_name=client.sparse_embedding_model_name,    batch_size=batch_size,
    parallel=None,))


# print(text_bm25_embeddings)


# client.create_collection(
#     "third Try",
#     vectors_config={
#         "all-MiniLM-L6-v2":models.VectorParams(
#             size=len(text_model_embeddings[0]),
#             distance=models.Distance.COSINE,
#         ),
#         "image-vectors": models.VectorParams(
#             size=len(image_embeddings[0]),  # Размер плотного вектора для изображений
#             distance=models.Distance.COSINE,
#         )
#     },
#     sparse_vectors_config = client.get_fastembed_sparse_vector_params()
# )


def create_points_for_qdrant(dataset, text_model_embeddings, image_embeddings):
    """Генерує точки для Qdrant на основі датасета і ембеддингів."""
    points = []
    id_counter = 1  # Начальный идентификатор для точек

    def get_next_id():
        nonlocal id_counter
        id_val = id_counter
        id_counter += 1
        return id_val
    
    i = 0

    for _,presentation in enumerate(dataset):
        for _,slide in enumerate(presentation["slides"]):
            # Додаємо текстову точку, якщо є текст
            points.append(
                    PointStruct(
                        id=get_next_id(), 
                        vector={
                            "all-MiniLM-L6-v2": text_model_embeddings[i],
                            client.get_sparse_vector_field_name(): encoded_sparse_docs[i],
                        },
                        payload={
                            "type": "text",
                            "file_name": presentation["file_name"],
                            "slide_number": slide["slide_number"],
                            "text": slide["text"],
                        },
                    )
                )
            i += 1
            for image_index, image_path in enumerate(slide["image_paths"]):
                points.append(
                    PointStruct(
                        id=get_next_id(),  # Використання UUID
                        vector={
                            "image-vectors": image_embeddings.pop(0),  # Беремо перший ембеддинг з черги
                        },
                        payload={
                            "type": "image",
                            "file_name": presentation["file_name"],
                            "slide_number": slide["slide_number"],
                            "image_path": image_path,
                        },
                    )
                )
    return points


# Генерація точок
points = create_points_for_qdrant(dataset, text_model_embeddings, image_embeddings)

# Додавання точок у Qdrant
client.upload_points(
    collection_name="third Try",
    points=points
)

# query_text = "What is the most popular Game engine?"

# dense_vectors, image_vectors=[],[]
# for query in tqdm.tqdm(query_text):
#     dense_query_vector=next(text_model.query_embed(texts))
#     # sparse_query_vector= next(text_bm25.query_embed(texts))
#     image_query_vector= next(image_model.embed(images))

#     dense_vectors.append(dense_query_vector)
#     # sparse_vectors.append(sparse_query_vector)
#     image_vectors.append(image_query_vector)

# print(dense_vectors)


# # Пример выполнения запроса и вывода результата
# response = client.query_points(
#     'third Try',
#     query=next(text_model.query_embed(query_text)),
#     using="all-MiniLM-L6-v2",
#     limit=10,
#     with_payload=False,
# )

# # Выводим результат запроса
# print(response)